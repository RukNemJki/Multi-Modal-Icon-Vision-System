from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_points):
    slide_layout = prs.slide_layouts[1]  # Bullet point layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title configuration
    title = slide.shapes.title
    title.text = title_text
    
    # Content configuration
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()  # Clear default empty paragraph
    
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)

# Initialize Presentation
prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]
title.text = "Multi-Modal Icon Vision System\nDeep Learning for Mobile UI Understanding"
subtitle.text = "Capstone Project Report (CPG-297)\n\nTeam: Harshit Sharma, Sushant Thakur, Kamal\nSupervisor: Dr. Surjit Singh\nThapar Institute of Engineering & Technology"

# Slide 2: Agenda
add_slide(prs, "Agenda", [
    "Problem Statement & Motivation",
    "Objectives & Scope",
    "Literature Review",
    "Proposed Solution: LYNX",
    "Methodology & Dataset",
    "System Architecture",
    "Implementation Details",
    "Results & Performance Analysis",
    "Real-World Applications",
    "Conclusion & Future Scope"
])

# Slide 3: Problem Statement
add_slide(prs, "Problem Statement", [
    "The Gap: Computers lack 'Native GUI Understanding'—cannot perceive screens like humans.",
    "Current Limitations:",
    "- Metadata (view hierarchies) often broken or inaccessible.",
    "- Existing detectors too slow for real-time interaction.",
    "Critical Needs:",
    "- Accessibility: Blind users need accurate alt-text for icons.",
    "- Automation: QA tools need visual grounding to test apps."
])

# Slide 4: Project Objectives
add_slide(prs, "Project Objectives", [
    "1. Dataset: Curate and preprocess Rico dataset for visual grounding.",
    "2. Model: Develop LYNX (Low-latency YOLO Nano eXtractor).",
    "3. Fusion: Integrate OCR for multi-modal reasoning.",
    "4. Deployment: Create containerized, production-ready API.",
    "5. Evaluation: Achieve high mAP and FPS for embedded use."
])

# Slide 5: Literature Review
add_slide(prs, "Literature Review", [
    "Evolution of Detection:",
    "- Two-Stage: Faster R-CNN is accurate but slow (~5 FPS).",
    "- One-Stage: YOLO allows real-time processing.",
    "Why YOLOv11?",
    "- Best trade-off between speed and accuracy.",
    "- ~30% faster training than YOLOv8.",
    "Gap: Most research focuses on layout, not pixel-level icon grounding[cite: 146]."
])

# Slide 6: Proposed Solution - LYNX
add_slide(prs, "Proposed Solution: LYNX", [
    "System Name: LYNX (Low-latency YOLO Nano eXtractor).",
    "Innovation: Multi-modal system fusing Vision (CV) and OCR.",
    "Key Features:",
    "- Visual Grounding: Locates icons in precise coordinates.",
    "- Semantic Fusion: Correlates icons with nearby text.",
    "- Extreme Speed: Designed for <1ms latency per image."
])

# Slide 7: Dataset Preparation
add_slide(prs, "Dataset Preparation", [
    "Source: Rico Dataset (72,000+ screenshots).",
    "Preprocessing:",
    "- Mapped raw data to 26 standardized icon classes.",
    "- Normalized bounding boxes to [0,1] format.",
    "Augmentation:",
    "- Mosaic & MixUp: Forces context learning.",
    "- HSV Augmentation: Handles dark/light themes."
])

# Slide 8: Model Architecture
add_slide(prs, "Model Architecture (YOLOv11 Nano)", [
    "Backbone (CSPDarknet + C3k2): Extracts multi-scale features for small icons.",
    "Neck (SPPF + PAN):",
    "- SPPF expands receptive field.",
    "- PAN fuses features for robust localization.",
    "Head: Anchor-free detection using Task-Aligned Learning (TAL)[cite: 370]."
])

# Slide 9: Multi-Modal Fusion Algorithm
add_slide(prs, "Multi-Modal Fusion Algorithm", [
    "Parallel Processing: YOLO and EasyOCR run simultaneously.",
    "Fusion Logic (Algorithm 1):",
    "- Calculates Euclidean distance between icon and text.",
    "- Proximity Threshold: Text within 100 pixels is associated.",
    "- Semantic Scoring: Matches keywords to validate predictions."
])

# Slide 10: System Architecture
add_slide(prs, "System Architecture", [
    "Input: Resize images to 640x640 and normalize.",
    "Four-Phase Pipeline:",
    "1. Ingestion: Validation and privacy filtering.",
    "2. Inference: LYNX (Vision) + EasyOCR (Text).",
    "3. Fusion: Spatial correlation of outputs.",
    "4. Output: JSON response with confidence scores."
])

# Slide 11: Technology Stack
add_slide(prs, "Technology Stack", [
    "Deep Learning: PyTorch 2.1, Ultralytics YOLOv11.",
    "OCR Engine: EasyOCR (Lightweight, multilingual).",
    "Backend: Flask API with Redis queuing[cite: 423].",
    "Deployment: Docker + TensorRT acceleration.",
    "Monitoring: Prometheus + Grafana."
])

# Slide 12: Sequence Diagram (API Flow)
add_slide(prs, "Sequence Diagram (API Flow)", [
    "1. Client posts image to /predict endpoint.",
    "2. System validates MIME type and JWT token.",
    "3. Parallel dispatch to Detection and OCR modules.",
    "4. Fusion module combines signals.",
    "5. Returns JSON response in <50ms."
])

# Slide 13: Implementation & DevOps
add_slide(prs, "Implementation & DevOps", [
    "Hugging Face Hosting:",
    "- Model weights and demo are public.",
    "- Enables zero-setup evaluation.",
    "Optimization:",
    "- Mixed Precision (AMP) & TF32 sped up training by 47%.",
    "- Exported to ONNX and TensorRT.",
    "QA: Nightly regression tests on 10k screenshots."
])

# Slide 14: Results - Overall Performance
add_slide(prs, "Results: Overall Performance", [
    "mAP50: 43.5% (Exceeds 40% target).",
    "Inference Speed: 1111 FPS (0.9ms per image).",
    "Precision: 52.3% (Conservative prediction).",
    "Latency: P99 is 67ms (Within 200ms limit)."
])

# Slide 15: Training Dynamics
add_slide(prs, "Training Dynamics", [
    "Loss Convergence: Stable learning without oscillation.",
    "mAP Trajectory: Reached 40% by epoch 30, plateaued at 43.5%.",
    "Efficiency: Training took ~4.4 hours on NVIDIA RTX 3060."
])

# Slide 16: Comparative Analysis
add_slide(prs, "Comparative Analysis", [
    "Faster R-CNN: 42.1% mAP | 5 FPS | 41.4M Params",
    "YOLOv8n: 37.3% mAP | 833 FPS | 3.2M Params",
    "OmniParser: 38.2% mAP | 45 FPS | 85.0M Params",
    "LYNX (Ours): 43.5% mAP | 1111 FPS | 2.6M Params",
    "Conclusion: LYNX is 6.2% more accurate & 33% faster than baseline [cite: 898-901]."
])

# Slide 17: Ablation Study
add_slide(prs, "Ablation Study", [
    "Baseline (YOLOv8n): 37.3% accuracy.",
    "+ YOLOv11 Arch: Improved to 41.5%.",
    "+ AMP/TF32: Reduced time by 47% (No accuracy loss).",
    "+ Full Augmentation: Final boost to 43.5% mAP50.",
    "Takeaway: Architecture + Augmentation drive gains."
])

# Slide 18: Quality of Service (QoS)
add_slide(prs, "Quality of Service (QoS)", [
    "Load Testing: 500 concurrent users.",
    "Throughput: 1,847 requests/sec (Target: 500).",
    "Availability: 99.97% uptime.",
    "Latency Distribution: 73% requests complete within 20ms."
])

# Slide 19: Real-World Impact
add_slide(prs, "Real-World Impact", [
    "Accessibility Case Study:",
    "- Found 312 unlabeled icons in 4,200 screens.",
    "- Reduced manual audit time by 27%.",
    "QA Automation:",
    "- Integrated into nightly regression suites.",
    "Business Impact: Estimated $110,000 annual savings."
])

# Slide 20: Conclusion & Future Scope
add_slide(prs, "Conclusion & Future Scope", [
    "Conclusion: Built production system with 43.5% accuracy and 1111 FPS.",
    "Future Scope:",
    "- Transformers: Use ViT for better context.",
    "- Multilingual: Support non-English apps via OCR.",
    "- Layout Understanding: Detect lists, grids, and cards."
])

# Save the presentation
prs.save('MultiModal_Icon_Vision_Presentation.pptx')
print("Presentation generated successfully.")